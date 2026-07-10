import os
import sys

# Support UTF-8 console output sur Windows
if hasattr(sys.stdout, 'reconfigure'):
    try:
        sys.stdout.reconfigure(encoding='utf-8', errors='replace')
    except Exception:
        pass

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("[ERREUR] La bibliotheque 'python-pptx' n'est pas installee.")
    print("Installez-la avec : pip install python-pptx pillow")
    sys.exit(1)

try:
    from PIL import Image, ImageDraw, ImageFont
    HAS_PIL = True
except ImportError:
    HAS_PIL = False
    print("[ATTENTION] Pillow non installe. Les images de diagramme seront remplacees par des formes natives.")

# ==========================================
# PALETTE DE COULEURS EVERTRUCK (NAVY/CRIMSON)
# ==========================================
NAVY = RGBColor(0, 29, 61)       # #001d3d
NAVY_LIGHT = RGBColor(0, 57, 112) # #003970
CRIMSON = RGBColor(230, 57, 70)  # #e63946
TEAL = RGBColor(15, 163, 177)    # #0fa3b1
GOLDEN = RGBColor(255, 159, 28)  # #ff9f1c
WHITE = RGBColor(255, 255, 255)
GRAY_BG = RGBColor(245, 247, 250)
GRAY_TEXT = RGBColor(100, 110, 120)
BLACK = RGBColor(20, 25, 30)

def create_architecture_image(filename="architecture_schema.png"):
    """Génère une image PNG schématisant l'architecture Docker du Monorepo via Pillow."""
    if not HAS_PIL:
        return None
    width, height = 1000, 550
    img = Image.new("RGB", (width, height), color=(245, 247, 250))
    draw = ImageDraw.Draw(img)

    # Fonction pour dessiner un bloc avec bordure et texte
    def draw_box(box, bg_color, border_color, title, subtitle="", text_color=(255,255,255)):
        x1, y1, x2, y2 = box
        draw.rectangle([x1, y1, x2, y2], fill=bg_color, outline=border_color, width=3)
        # Titre
        draw.text(((x1+x2)//2, (y1+y2)//2 - (15 if subtitle else 0)), title, fill=text_color, anchor="mm")
        if subtitle:
            draw.text(((x1+x2)//2, (y1+y2)//2 + 15), subtitle, fill=(220,230,240) if text_color==(255,255,255) else (100,110,120), anchor="mm")

    # Conteneur principal Docker Network
    draw.rectangle([180, 40, 960, 510], outline=(0, 29, 61), width=2)
    draw.text((195, 55), "🐳 Docker Bridge Network : evertruck-network", fill=(0, 29, 61))

    # Blocs
    draw_box([20, 220, 150, 320], (255, 255, 255), (230, 57, 70), "🌐 Clients", "Web / Mobile", text_color=(0, 29, 61))
    draw_box([210, 210, 360, 330], (0, 29, 61), (230, 57, 70), "🛡️ Nginx Proxy", "Ports 80 / 443\nSSL & Gzip")
    
    draw_box([430, 110, 630, 230], (15, 163, 177), (255, 255, 255), "⚡ Frontend", "Port 3000\nReact 19 + Three.js")
    draw_box([430, 310, 630, 430], (230, 57, 70), (255, 255, 255), "🚀 Backend API", "Port 4000\nHono + tRPC 11")
    
    draw_box([720, 240, 920, 350], (255, 159, 28), (255, 255, 255), "💾 MySQL 8.0", "Port 3306\n14 Tables Drizzle", text_color=(0, 0, 0))
    draw_box([720, 380, 920, 480], (180, 30, 40), (255, 255, 255), "🔥 Redis 7", "Port 6379\nCache & Sessions")

    # Flèches / Lignes
    def draw_arrow(start, end, color=(100,110,120)):
        draw.line([start, end], fill=color, width=3)

    draw_arrow((150, 270), (210, 270), (230, 57, 70))
    draw_arrow((360, 250), (430, 170), (15, 163, 177))
    draw_arrow((360, 290), (430, 370), (230, 57, 70))
    draw_arrow((630, 350), (720, 295), (255, 159, 28))
    draw_arrow((630, 390), (720, 430), (180, 30, 40))

    img.save(filename)
    return filename

def create_kpi_image(filename="stats_chart.png"):
    """Génère un graphique visuel des métriques clés de la plateforme."""
    if not HAS_PIL:
        return None
    width, height = 900, 400
    img = Image.new("RGB", (width, height), color=(0, 29, 61))
    draw = ImageDraw.Draw(img)

    metrics = [
        ("5", "Conteneurs Docker\nIsolés & Sécurisés", (230, 57, 70)),
        ("14", "Tables Relationnelles\nMySQL 8 + Drizzle", (255, 159, 28)),
        ("15", "Routeurs tRPC 11\nTypés de bout-en-bout", (15, 163, 177)),
        ("0%", "Dette de Typage\nContrat API Garanti", (40, 200, 120))
    ]

    box_w = 190
    spacing = 30
    start_x = (width - (len(metrics) * box_w + (len(metrics)-1) * spacing)) // 2

    for i, (val, label, color) in enumerate(metrics):
        bx1 = start_x + i * (box_w + spacing)
        by1, bx2, by2 = 60, bx1 + box_w, 340
        draw.rectangle([bx1, by1, bx2, by2], fill=(0, 45, 90), outline=color, width=3)
        # Valeur
        draw.text(((bx1+bx2)//2, by1 + 80), val, fill=color, anchor="mm")
        # Label
        draw.text(((bx1+bx2)//2, by1 + 180), label, fill=(255,255,255), anchor="mm")

    img.save(filename)
    return filename

def generate_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Helper pour ajouter un header de slide
    def add_header(slide, title_text, category_text="AUDIT TECHNIQUE EVERTRUCK"):
        # Bandeau haut Navy
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = NAVY
        header_shape.line.color.rgb = NAVY

        # Accent Crimson
        accent_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(13.333), Inches(0.08))
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = CRIMSON
        accent_shape.line.color.rgb = CRIMSON

        # Texte Catégorie
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(10), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(11)
        p.font.color.rgb = TEAL
        p.font.bold = True

        # Texte Titre
        txBox2 = slide.shapes.add_textbox(In.ches(0.6) if hasattr(Inches, 'ches') else Inches(0.6), Inches(0.38), Inches(12), Inches(0.6))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(26)
        p2.font.color.rgb = WHITE
        p2.font.bold = True

    # ──────────────────────────────────────────
    # SLIDE 1 : TITRE / HERO
    # ──────────────────────────────────────────
    slide_layout = prs.slide_layouts[6] # Blanc vide
    slide1 = prs.slides.add_slide(slide_layout)
    
    # Fond marine global
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.color.rgb = NAVY

    # Carré décoratif rouge
    deco = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.8), Inches(0.2), Inches(3.2))
    deco.fill.solid()
    deco.fill.fore_color.rgb = CRIMSON
    deco.line.color.rgb = CRIMSON

    txBox = slide1.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.5), Inches(2.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "EVERTRUCK LOGISTICS"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = "Audit Architecture & Rapport Technique — Production-Ready"
    p2.font.size = Pt(24)
    p2.font.color.rgb = TEAL

    p3 = tf.add_paragraph()
    p3.text = "\nPlateforme Monorepo TypeScript • React 19 • Three.js • Hono • tRPC 11 • MySQL 8 • Docker"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(180, 200, 220)

    # ──────────────────────────────────────────
    # SLIDE 2 : RÉSUMÉ EXÉCUTIF & CHIFFRES CLÉS
    # ──────────────────────────────────────────
    slide2 = prs.slides.add_slide(slide_layout)
    add_header(slide2, "Résumé Exécutif & Chiffres Clés")

    # Image KPI
    kpi_img = create_kpi_image()
    if kpi_img and os.path.exists(kpi_img):
        slide2.shapes.add_picture(kpi_img, Inches(0.6), Inches(1.5), width=Inches(12.1))
    else:
        # Fallback texte si pas d'image
        tx = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
        tx.text_frame.text = "5 Conteneurs Docker | 14 Tables MySQL | 15 Routeurs tRPC | 0% Dette de Typage"

    tx_bot = slide2.shapes.add_textbox(Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.5))
    tf_bot = tx_bot.text_frame
    tf_bot.word_wrap = True
    p = tf_bot.paragraphs[0]
    p.text = "🎯 Objectif : Présenter la structure et les atouts d'EverTruck, une plateforme logistique conçue sous la forme d'un Monorepo moderne et hautement sécurisé pour la gestion de fret en temps réel."
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

    # ──────────────────────────────────────────
    # SLIDE 3 : ARCHITECTURE GLOBALE DOCKER
    # ──────────────────────────────────────────
    slide3 = prs.slides.add_slide(slide_layout)
    add_header(slide3, "Architecture Monorepo & Conteneurisation")

    arch_img = create_architecture_image()
    if arch_img and os.path.exists(arch_img):
        slide3.shapes.add_picture(arch_img, Inches(1.2), Inches(1.5), width=Inches(10.9))

    tx3 = slide3.shapes.add_textbox(Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.2))
    tf3 = tx3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "🐳 Isolation complète réseau (evertruck-network) : Nginx gère le SSL et le Gzip en amont, le Frontend sert l'UI et les éléments 3D (Port 3000), le Backend Hono/tRPC pilote MySQL 8 et Redis 7 (Port 4000)."
    p3.font.size = Pt(16)
    p3.font.color.rgb = NAVY

    # ──────────────────────────────────────────
    # SLIDE 4 : AUDIT FRONTEND (REACT 19 + 3D)
    # ──────────────────────────────────────────
    slide4 = prs.slides.add_slide(slide_layout)
    add_header(slide4, "Audit Frontend : React 19, UI/UX Premium & Rendu 3D", "COUCHE CLIENT / SPA")

    # 2 colonnes
    col1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.4))
    col1.fill.solid()
    col1.fill.fore_color.rgb = GRAY_BG
    col1.line.color.rgb = TEAL
    
    tx_col1 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.4), Inches(5.0))
    tf_c1 = tx_col1.text_frame
    tf_c1.word_wrap = True
    p = tf_c1.paragraphs[0]
    p.text = "⚡ Stack Technologique & Design"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = NAVY

    points_c1 = [
        "• React 19 + TypeScript + Vite 7 pour un build instantané et une fluidité maximale.",
        "• Palette de couleurs exclusive : Navy (#001d3d) institutionnel, Crimson (#e63946) pour les alertes/CTA, et Teal (#0fa3b1).",
        "• Immersion 3D (Three.js / React Three Fiber) : Animations et scènes interactives sur la page d'accueil.",
        "• Animations Framer Motion & GSAP avec défilement doux (Lenis Scroll)."
    ]
    for pt in points_c1:
        p = tf_c1.add_paragraph()
        p.text = pt
        p.font.size = Pt(15)
        p.font.color.rgb = BLACK
        p.space_before = Pt(12)

    col2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.4))
    col2.fill.solid()
    col2.fill.fore_color.rgb = GRAY_BG
    col2.line.color.rgb = CRIMSON

    tx_col2 = slide4.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.4), Inches(5.0))
    tf_c2 = tx_col2.text_frame
    tf_c2.word_wrap = True
    p = tf_c2.paragraphs[0]
    p.text = "📑 Structure & Performance des Pages"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = CRIMSON

    points_c2 = [
        "• Chargement Eager / Lazy : Les pages critiques sont instantanées, les pages secondaires déportées via React.lazy().",
        "• Pages Vitrine : Accueil avec compteurs en temps réel, Flotte interactive, Demande de devis multi-étapes avec calcul de volume.",
        "• Espace Client / Tracking : Suivi en direct du statut d'acheminement (pending, in_transit, delivered).",
        "• Dashboard Admin : 11 onglets pour piloter l'intégralité du CRM et du parc routier sans requête SQL manuelle."
    ]
    for pt in points_c2:
        p = tf_c2.add_paragraph()
        p.text = pt
        p.font.size = Pt(15)
        p.font.color.rgb = BLACK
        p.space_before = Pt(12)

    # ──────────────────────────────────────────
    # SLIDE 5 : AUDIT BACKEND (HONO + tRPC)
    # ──────────────────────────────────────────
    slide5 = prs.slides.add_slide(slide_layout)
    add_header(slide5, "Audit Backend : Hono, tRPC 11 & Sécurité Zod", "COUCHE SERVEUR API")

    tx5 = slide5.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.3))
    tf5 = tx5.text_frame
    tf5.word_wrap = True

    items5 = [
        ("🚀 Moteur Hono Ultra-Rapide", "Utilisé pour sa légèreté extrême et sa vitesse d'exécution sous Node 20. Intègre une limite de payload généreuse de 50 MB pour l'upload de photos de fret et de documents de transport (BOL)."),
        ("🔗 Routeur tRPC 11 (15 Sous-Modules)", "Sépare parfaitement la logique métier en routeurs spécialisés : auth, localAuth, vehicle, shipment, tracking, quote, invoice, blog, faq, testimonial, contact, partner, notification, user, stats."),
        ("🛡️ Validation Zod & Contrat de Typage", "Chaque payload d'entrée (formulaire de devis, création d'expédition, connexion) est rigoureusement validé par Zod en amont. Zéro risque d'injection ou d'erreur d'interface."),
        ("🔐 Authentification & Sécurité JWT", "Système double d'authentification (OAuth Kimi + Local Auth email/password Bcrypt). Les tokens JWT sont signés de manière cryptographique et transmis via des cookies httpOnly et SameSite sécurisés.")
    ]

    for i, (title, desc) in enumerate(items5):
        p = tf5.paragraphs[0] if i == 0 else tf5.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        if i > 0:
            p.space_before = Pt(14)
        
        p_desc = tf5.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(15)
        p_desc.font.color.rgb = BLACK
        p_desc.space_before = Pt(4)

    # ──────────────────────────────────────────
    # SLIDE 6 : BASE DE DONNÉES MYSQL & SCHEMA
    # ──────────────────────────────────────────
    slide6 = prs.slides.add_slide(slide_layout)
    add_header(slide6, "Modèle de Données : MySQL 8.0 & Drizzle ORM", "COUCHE PERSISTANCE")

    # Tableau des tables
    rows, cols = 8, 3
    table_shape = slide6.shapes.add_table(rows, cols, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.3))
    table = table_shape.table
    table.columns[0].width = Inches(2.3)
    table.columns[1].width = Inches(3.6)
    table.columns[2].width = Inches(6.2)

    headers = ["Table MySQL", "Rôle Métier", "Spécificités & Contraintes"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.size = Pt(14)

    data_rows = [
        ("users", "Comptes Clients & Admins", "Identifiants uniques (unionId), hachage Bcrypt, rôles (user, admin, manager)"),
        ("vehicles", "Flotte de Camions", "Types (truck, van, trailer, forklift, crane), immatriculation, capacité en tonnes"),
        ("drivers", "Chauffeurs Routiers", "Permis CDL (Class A/B), note moyenne (rating), statut (active, off_duty)"),
        ("shipments", "Expéditions de Fret", "Numéro de suivi EVT-XXXX, origine/destination, poids, statut de livraison"),
        ("tracking", "Historique GPS", "Lié aux shipments, horodatage, coordonnées (latitude, longitude) & description"),
        ("invoices", "Facturation & Taxes", "Montants H.T., TVA, Total TTC, statut (draft, sent, paid, overdue)"),
        ("quote_requests / contact", "Prospects & Support CRM", "Calcul de volume de fret, suivi d'examen des demandes, archivage de support")
    ]

    for r_idx, row_data in enumerate(data_rows):
        for c_idx, text_val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = text_val
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = BLACK
                if c_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = CRIMSON

    # ──────────────────────────────────────────
    # SLIDE 7 : DEVOPS, NGINX & SÉCURITÉ
    # ──────────────────────────────────────────
    slide7 = prs.slides.add_slide(slide_layout)
    add_header(slide7, "DevOps, Sécurité Nginx & Déploiement VPS", "INFRASTRUCTURE & PROD")

    tx7 = slide7.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.3))
    tf7 = tx7.text_frame
    tf7.word_wrap = True

    items7 = [
        ("🛡️ En-têtes de Sécurité Strictes (Nginx Alpine)", "Le reverse proxy applique des règles incontournables : X-Frame-Options SAMEORIGIN anti-clickjacking, X-XSS-Protection, X-Content-Type-Options nosniff, et Content-Security-Policy (CSP)."),
        ("⚡ Optimisation & Compression de Bande Passante", "Compression Gzip de niveau 6 active sur l'ensemble des flux JSON, CSS, JS et SVG. Les actifs statiques du frontend bénéficient d'une mise en cache longue durée d'1 an (immutable)."),
        ("🐳 Builds Multi-Stage Docker", "Les images de production (Dockerfile.prod) compilent le TypeScript et les assets Vite dans des étapes temporaires (builders), ne copiant que les artefacts finaux dans des images Alpine légères."),
        ("🔑 Sécurité des Variables d'Environnement (.env.production)", "Isolation stricte entre le développement local et la production. Les secrets cryptographiques et mots de passe bases de données sont externalisés via des templates sécurisés.")
    ]

    for i, (title, desc) in enumerate(items7):
        p = tf7.paragraphs[0] if i == 0 else tf7.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = CRIMSON
        if i > 0:
            p.space_before = Pt(14)
        
        p_desc = tf7.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(15)
        p_desc.font.color.rgb = BLACK
        p_desc.space_before = Pt(4)

    # ──────────────────────────────────────────
    # SLIDE 8 : FORCES & RECOMMANDATIONS
    # ──────────────────────────────────────────
    slide8 = prs.slides.add_slide(slide_layout)
    add_header(slide8, "Forces du Projet & Feuille de Route", "SYNTHÈSE & NEXT STEPS")

    b1 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2))
    b1.fill.solid()
    b1.fill.fore_color.rgb = RGBColor(235, 245, 250)
    b1.line.color.rgb = TEAL

    tx_b1 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_b1 = tx_b1.text_frame
    tf_b1.word_wrap = True
    p = tf_b1.paragraphs[0]
    p.text = "✅ Forces Majeures d'EverTruck"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = NAVY

    pts_b1 = [
        "• Zéro Dette Technique : Couplage parfait TypeScript + tRPC + Drizzle ORM.",
        "• UX/UI Haut de Gamme : Rendu 3D sur le Héro, compteurs interactifs et design soigné.",
        "• Autonomie Administrative : AdminDashboard de 11 onglets couvrant 100% du cycle de vie logistique.",
        "• Sécurité et Isolation : Réseau Docker privé avec proxy Nginx en bouclier."
    ]
    for pt in pts_b1:
        p = tf_b1.add_paragraph()
        p.text = pt
        p.font.size = Pt(15)
        p.font.color.rgb = BLACK
        p.space_before = Pt(12)

    b2 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2))
    b2.fill.solid()
    b2.fill.fore_color.rgb = RGBColor(255, 243, 243)
    b2.line.color.rgb = CRIMSON

    tx_b2 = slide8.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_b2 = tx_b2.text_frame
    tf_b2.word_wrap = True
    p = tf_b2.paragraphs[0]
    p.text = "🚀 Recommandations d'Évolution"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = CRIMSON

    pts_b2 = [
        "• Suivi GPS WebSockets / SSE : Animer en continu la position des camions sur une carte interactive sans rechargement.",
        "• Export PDF des Factures : Intégrer Puppeteer / PDFKit côté serveur pour télécharger instantanément les factures et bons de livraison.",
        "• Notifications SMS & Email : Relier notificationRouter à SendGrid ou Twilio pour prévenir le client à chaque étape d'acheminement."
    ]
    for pt in pts_b2:
        p = tf_b2.add_paragraph()
        p.text = pt
        p.font.size = Pt(15)
        p.font.color.rgb = BLACK
        p.space_before = Pt(14)

    # Sauvegarde avec gestion de verrouillage (si le fichier est ouvert dans PowerPoint)
    output_pptx = "EverTruck_Presentation_Audit.pptx"
    base_name, ext = os.path.splitext(output_pptx)
    version = 1
    while True:
        try:
            prs.save(output_pptx)
            print(f"[SUCCES] Presentation PowerPoint creee avec succes : {output_pptx}")
            break
        except PermissionError:
            version += 1
            output_pptx = f"{base_name}_v{version}{ext}"
            print(f"[INFO] Fichier verrouille, tentative de sauvegarde sous : {output_pptx}")

if __name__ == "__main__":
    generate_presentation()
